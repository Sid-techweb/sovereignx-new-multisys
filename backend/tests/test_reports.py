import io
import unittest
from fastapi.testclient import TestClient
from app.main import app
from app.config import settings
import docx

class TestReportsApi(unittest.TestCase):
    def setUp(self):
        self.client = TestClient(app, headers={"X-API-Key": settings.API_KEY})
        
        # Build mock AgentInvestigateResponse input payload
        self.mock_payload = {
            "query": "What is the status of Pump P-204?",
            "answer": (
                "Pump P-204 experienced abnormal housing temperatures of 91 C "
                "[Source: pump_P204_sensor_data.csv | chunk_id=c2c4c44f-4bbb-420d-898d-449ed40a9f02]. "
                "The radial vibration reading was elevated at 5.8 mm/s "
                "[Source: pump_P204_sensor_data.csv | chunk_id=c2c4c44f-4bbb-420d-898d-449ed40a9f02]. "
                "The SOP bearing housing limit is 80 C "
                "[Source: pump_P204_SOP.pdf | page=1 | chunk_id=377c635a-2a55-4de3-b040-522c4bb00973]."
            ),
            "retrieved_chunks": [
                {
                    "chunk_id": "c2c4c44f-4bbb-420d-898d-449ed40a9f02",
                    "document_id": "doc-001",
                    "filename": "pump_P204_sensor_data.csv",
                    "source": "user_upload",
                    "content": "temperature_c: 91\nvibration_mm_s: 5.8",
                    "score": 0.85,
                    "metadata": {
                        "page_number": None,
                        "chunk_index": 0
                    }
                },
                {
                    "chunk_id": "377c635a-2a55-4de3-b040-522c4bb00973",
                    "document_id": "doc-002",
                    "filename": "pump_P204_SOP.pdf",
                    "source": "user_upload",
                    "content": "bearing temperature limit is 80 C",
                    "score": 0.80,
                    "metadata": {
                        "page_number": 1,
                        "chunk_index": 0
                    }
                }
            ],
            "confidence": 0.825,
            "tool_executions": [
                {
                    "tool_name": "compare_reading_against_sop_limit",
                    "arguments": {
                        "reading_value": 91.0,
                        "limit_value": 80.0,
                        "comparison_type": "greater_than",
                        "unit": "C"
                    },
                    "status": "success",
                    "outputs": {
                        "is_exceeded": True,
                        "summary": "Exceedance detected: 91.0 C is greater than limit of 80.0 C"
                    },
                    "context_id": "test-context"
                }
            ],
            "metadata": {
                "model_used": "qwen2.5:7b",
                "latency_ms": 1250.0
            }
        }

    def test_generate_docx_success(self):
        """Test POST /reports/generate-docx responds with valid Word docx and correct mappings"""
        response = self.client.post("/reports/generate-docx", json=self.mock_payload)
        self.assertEqual(response.status_code, 200)
        
        # Verify content type
        content_type = response.headers.get("content-type")
        self.assertEqual(content_type, "application/vnd.openxmlformats-officedocument.wordprocessingml.document")
        
        # Load returned bytes into python-docx Document to verify validity
        doc_bytes = io.BytesIO(response.content)
        doc = docx.Document(doc_bytes)
        
        # Extract all text paragraphs
        paragraphs_text = [p.text for p in doc.paragraphs]
        full_text = "\n".join(paragraphs_text)
        
        # 1. Assert Title is present
        self.assertIn("SovereignX Grounded Investigation Report", full_text)
        
        # 2. Assert query and mapped answer with bracketed citation markers are present
        self.assertIn('What is the status of Pump P-204?', full_text)
        
        # The answer contains 2 unique citations. The CSV is mentioned twice, so it should get [1].
        # The SOP should get [2].
        expected_answer_text = (
            "Pump P-204 experienced abnormal housing temperatures of 91 C [1]. "
            "The radial vibration reading was elevated at 5.8 mm/s [1]. "
            "The SOP bearing housing limit is 80 C [2]."
        )
        self.assertIn(expected_answer_text, full_text)
        
        # 3. Verify tool execution results are rendered and distinct
        self.assertIn("Deterministic Tool Verifications (Phase 6)", full_text)
        self.assertIn("Exceedance detected: 91.0 C is greater than limit of 80.0 C", [c.text for row in doc.tables[0].rows for c in row.cells])
        
        # 4. Verify confidence score description is correct
        self.assertIn("Average Retrieval Relevance Confidence: 82.5%", full_text)
        self.assertIn("This score represents the mathematical average", full_text)
        
        # 5. Verify unique citations inside the sources table (second table in the doc)
        sources_table = doc.tables[1]
        rows_content = [[cell.text for cell in row.cells] for row in sources_table.rows]
        
        # Assert de-duplicated rows (1 header + 2 unique citation rows = 3 rows total)
        self.assertEqual(len(rows_content), 3)
        
        # Assert Ref [1] matches CSV info
        self.assertEqual(rows_content[1][0], "[1]")
        self.assertEqual(rows_content[1][1], "pump_P204_sensor_data.csv")
        self.assertEqual(rows_content[1][2], "N/A")
        self.assertEqual(rows_content[1][3], "c2c4c44f-4bbb-420d-898d-449ed40a9f02")
        
        # Assert Ref [2] matches SOP info
        self.assertEqual(rows_content[2][0], "[2]")
        self.assertEqual(rows_content[2][1], "pump_P204_SOP.pdf")
        self.assertEqual(rows_content[2][2], "1")
        self.assertEqual(rows_content[2][3], "377c635a-2a55-4de3-b040-522c4bb00973")

    def test_generate_pptx_success(self):
        """Test POST /reports/generate-pptx responds with valid PowerPoint presentation"""
        response = self.client.post("/reports/generate-pptx", json=self.mock_payload)
        self.assertEqual(response.status_code, 200)
        
        # Verify content type
        content_type = response.headers.get("content-type")
        self.assertEqual(content_type, "application/vnd.openxmlformats-officedocument.presentationml.presentation")
        
        # Load returned bytes into python-pptx Presentation
        from pptx import Presentation
        prs = Presentation(io.BytesIO(response.content))
        
        # We expect 3 slides (Title, Findings, Sources)
        self.assertEqual(len(prs.slides), 3)
        
        # Check title slide content
        slide1 = prs.slides[0]
        slide1_text = "".join([shape.text for shape in slide1.shapes if shape.has_text_frame])
        self.assertIn("What is the status of Pump P-204?", slide1_text)
        self.assertIn("SovereignX Grounded Investigation Report", slide1_text)
        
        # Check findings slide content
        slide2 = prs.slides[1]
        slide2_text = "".join([shape.text for shape in slide2.shapes if shape.has_text_frame])
        self.assertIn("Grounded Findings & Analysis", slide2_text)
        self.assertIn("Pump P-204 experienced abnormal housing temperatures of 91 C [1].", slide2_text)
        self.assertIn("Confidence: 82.5%", slide2_text)

    def test_generate_xlsx_success(self):
        """Test POST /reports/generate-xlsx responds with valid Excel sheet and claims grid"""
        response = self.client.post("/reports/generate-xlsx", json=self.mock_payload)
        self.assertEqual(response.status_code, 200)
        
        # Verify content type
        content_type = response.headers.get("content-type")
        self.assertEqual(content_type, "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
        
        # Load returned bytes into openpyxl Workbook
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(response.content))
        
        # Check sheets
        self.assertIn("Provenance Audit", wb.sheetnames)
        self.assertIn("Tool Verifications", wb.sheetnames)
        
        ws = wb["Provenance Audit"]
        
        # Cell A1 should be title
        self.assertEqual(ws["A1"].value, "SovereignX Grounded Findings Claim-by-Claim Audit")
        
        # Row 5 should have headers
        self.assertEqual(ws.cell(row=5, column=1).value, "Claim Text / Grounded Finding")
        self.assertEqual(ws.cell(row=5, column=2).value, "Citation Ref")
        
        # Row 6 should contain the first sentence / claim
        claim_val = ws.cell(row=6, column=1).value
        ref_val = ws.cell(row=6, column=2).value
        file_val = ws.cell(row=6, column=3).value
        
        self.assertIn("Pump P-204 experienced abnormal housing temperatures of 91 C [1].", claim_val)
        self.assertEqual(ref_val, "[1]")
        self.assertEqual(file_val, "pump_P204_sensor_data.csv")
