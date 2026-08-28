import io
import re
from fastapi import APIRouter, HTTPException
from fastapi.responses import StreamingResponse
from app.schemas import AgentInvestigateResponse
import docx
from docx.shared import Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from typing import List, Dict, Any, Tuple

router = APIRouter(prefix="/reports", tags=["reports"])

def extract_unique_citations(raw_answer: str, retrieved_chunks: List[Dict[str, Any]]) -> Tuple[str, List[Dict[str, Any]]]:
    """
    Utility helper to parse [Source: ...] citation blocks out of raw answer text,
    de-duplicate them, assign sequential Ref # indices [1], [2], etc.,
    and map them back to metadata from retrieved chunks.
    """
    pattern = re.compile(r'\[Source:\s*([^\]]+)\]')
    matches = pattern.findall(raw_answer)
    
    unique_citations = []
    citation_map = {}  # raw_text -> index
    chunks_by_id = {c.get("chunk_id"): c for c in retrieved_chunks if c.get("chunk_id")}
    
    for match in matches:
        raw_text = f"[Source: {match}]"
        if raw_text in citation_map:
            continue
            
        parts = [p.strip() for p in match.split('|')]
        filename = parts[0] if parts else "unknown_file"
        
        chunk_id = None
        page = None
        
        for part in parts[1:]:
            if '=' in part:
                k, v = part.split('=', 1)
                k, v = k.strip(), v.strip()
                if k == 'chunk_id':
                    chunk_id = v
                elif k == 'page':
                    page = v
            elif part.startswith('page='):
                page = part.split('=', 1)[1].strip()
            elif part.startswith('chunk_id='):
                chunk_id = part.split('=', 1)[1].strip()
        
        # Correlate chunk details from RAG payload if available
        if chunk_id and chunk_id in chunks_by_id:
            chunk = chunks_by_id[chunk_id]
            filename = chunk.get("filename", filename)
            page = chunk.get("metadata", {}).get("page_number", page)
            
        index = len(unique_citations) + 1
        citation_map[raw_text] = index
        unique_citations.append({
            "index": index,
            "filename": filename,
            "page": page,
            "chunk_id": chunk_id or "unknown"
        })
        
    formatted_answer = raw_answer
    for raw_text, index in citation_map.items():
        formatted_answer = formatted_answer.replace(raw_text, f"[{index}]")
        
    return formatted_answer, unique_citations

def extract_inputs_from_summary(summary: str) -> str:
    """
    Parses reading and SOP limit values from a tool output summary string.
    """
    match = re.search(r'Reading \(([^)]+)\).*?SOP limit \(([^)]+)\)', summary, re.IGNORECASE)
    if match:
        return f"reading={match.group(1)}, limit={match.group(2)}"
    return "N/A"

@router.post("/generate-docx")
async def generate_docx(request: AgentInvestigateResponse):
    """
    Consumes an existing /agents/investigate response directly and returns a downloadable .docx file.
    """
    try:
        # Create in-memory document
        doc = docx.Document()
        
        # Style setup
        style = doc.styles['Normal']
        font = style.font
        font.name = 'Calibri'
        font.size = Pt(11)
        
        # 1. Document Title
        title = doc.add_paragraph()
        title_run = title.add_run("SovereignX Grounded Investigation Report")
        title_run.bold = True
        title_run.font.size = Pt(20)
        title_run.font.color.rgb = RGBColor(16, 44, 87)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # 2. Metadata / Summary Section
        doc.add_heading("Investigation Summary", level=2)
        
        p_meta = doc.add_paragraph()
        p_meta.add_run("User Query: ").bold = True
        p_meta.add_run(f"\"{request.query}\"\n")
        
        # Confidence Score with explicit label
        p_meta.add_run("Average Retrieval Relevance Confidence: ").bold = True
        p_meta.add_run(f"{round(request.confidence * 100, 2)}%\n")
        
        # Explanatory subtext for what confidence score represents
        metric_explanation = (
            "  (Metric description: This score represents the mathematical average of the vector "
            "similarity relevance scores of the retrieved document chunks used to ground this report.)\n\n"
        )
        desc_run = p_meta.add_run(metric_explanation)
        desc_run.font.size = Pt(9.5)
        desc_run.font.color.rgb = RGBColor(120, 120, 120)
        
        p_meta.add_run("Gateway Model Used: ").bold = True
        p_meta.add_run(f"{request.metadata.get('model_used', 'N/A')}\n")
        p_meta.add_run("Latency: ").bold = True
        p_meta.add_run(f"{request.metadata.get('latency_ms', 0.0):.2f} ms\n")
        
        # 3. Grounded Answer Section
        doc.add_heading("Grounded Findings & Analysis", level=2)
        
        formatted_answer, unique_citations = extract_unique_citations(request.answer, request.retrieved_chunks)
        doc.add_paragraph(formatted_answer)
        
        # 4. Deterministic Tool Executions Section (Phase 6)
        doc.add_heading("Deterministic Tool Verifications (Phase 6)", level=2)
        
        note_p = doc.add_paragraph()
        note_run = note_p.add_run(
            "Note: The following results are computed deterministically by SovereignX's comparison tool rules, completely independent of the LLM."
        )
        note_run.italic = True
        note_run.font.color.rgb = RGBColor(120, 120, 120)
        
        if not request.tool_executions:
            doc.add_paragraph("No comparison tools were triggered for this query context.")
        else:
            table = doc.add_table(rows=1, cols=4)
            table.style = 'Light Shading Accent 1'
            hdr_cells = table.rows[0].cells
            hdr_cells[0].text = 'Tool Name'
            hdr_cells[1].text = 'Parameters / Inputs'
            hdr_cells[2].text = 'Status'
            hdr_cells[3].text = 'Deterministic Output Summary'
            
            for t in request.tool_executions:
                row_cells = table.add_row().cells
                row_cells[0].text = str(t.get("tool_name", "Unknown"))
                
                summary = t.get("outputs", {}).get("summary", "No output summary")
                inputs_str = extract_inputs_from_summary(summary)
                if inputs_str == "N/A":
                    args = t.get("arguments", {})
                    if args:
                        inputs_str = ", ".join([f"{k}={v}" for k, v in args.items()])
                
                row_cells[1].text = inputs_str
                row_cells[2].text = str(t.get("status", "success"))
                row_cells[3].text = summary
                
        # 5. Sources Section
        doc.add_heading("Sources & Provenance Details", level=2)
        
        if not unique_citations:
            doc.add_paragraph("No evidence citations were referenced in the grounded findings.")
        else:
            table = doc.add_table(rows=1, cols=4)
            table.style = 'Medium Shading 1 Accent 1'
            hdr_cells = table.rows[0].cells
            hdr_cells[0].text = 'Ref #'
            hdr_cells[1].text = 'Document Filename'
            hdr_cells[2].text = 'Page #'
            hdr_cells[3].text = 'Source Chunk ID (Provenance)'
            
            for c in unique_citations:
                row_cells = table.add_row().cells
                row_cells[0].text = f"[{c['index']}]"
                row_cells[1].text = str(c['filename'])
                row_cells[2].text = str(c['page']) if c['page'] is not None else "N/A"
                row_cells[3].text = str(c['chunk_id'])
                
        # Return StreamingResponse of Docx
        file_stream = io.BytesIO()
        doc.save(file_stream)
        file_stream.seek(0)
        
        headers = {
            "Content-Disposition": "attachment; filename=SovereignX_Report.docx"
        }
        return StreamingResponse(
            file_stream,
            media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            headers=headers
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to generate Word document: {str(e)}")

@router.post("/generate-pptx")
async def generate_pptx(request: AgentInvestigateResponse):
    """
    Consumes an existing /agents/investigate response directly and returns a downloadable .pptx file.
    """
    try:
        from pptx import Presentation
        from pptx.util import Pt
        from pptx.dml.color import RGBColor as PptRGBColor
        
        prs = Presentation()
        
        # Helper to set text formatting
        def set_font(run, size_pt=14, bold=False, italic=False, color_rgb=None):
            run.font.name = "Calibri"
            run.font.size = Pt(size_pt)
            run.font.bold = bold
            run.font.italic = italic
            if color_rgb:
                run.font.color.rgb = color_rgb
        
        # Format answer and parse citations
        formatted_answer, unique_citations = extract_unique_citations(request.answer, request.retrieved_chunks)
        
        # 1. Slide 1: Title Slide (layout index 0 is Title Slide)
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = request.query
        subtitle.text = "SovereignX Grounded Investigation Report"
        
        # 2. Slide 2: Findings Slide (layout index 1 is Title and Content)
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        title_shape.text = "Grounded Findings & Analysis"
        
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.word_wrap = True
        
        # Split findings text into paragraphs for presentation
        paragraphs = [p.strip() for p in formatted_answer.split('\n') if p.strip()]
        if not paragraphs:
            paragraphs = ["No grounded findings could be established."]
            
        for i, p_text in enumerate(paragraphs):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = p_text
            p.font.size = Pt(14)
            
        # Add a confidence and model tag at the bottom
        p_conf = tf.add_paragraph()
        p_conf.space_before = Pt(20)
        run_conf = p_conf.add_run()
        run_conf.text = f"Confidence: {round(request.confidence * 100, 2)}% | Model: {request.metadata.get('model_used', 'N/A')}"
        set_font(run_conf, size_pt=12, italic=True, color_rgb=PptRGBColor(120, 120, 120))
        
        # 3. Slide 3: Sources Slide
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        title_shape.text = "Sources & Provenance Details"
        
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.word_wrap = True
        
        if not unique_citations:
            p = tf.paragraphs[0]
            p.text = "No evidence citations were referenced in the grounded findings."
        else:
            for i, c in enumerate(unique_citations):
                p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
                page_str = f", Page {c['page']}" if c['page'] is not None else ""
                p.text = f"[{c['index']}] {c['filename']}{page_str} (Chunk: {c['chunk_id'][:8]}...)"
                p.level = 0
                p.font.size = Pt(14)
                
        # Save to memory stream
        file_stream = io.BytesIO()
        prs.save(file_stream)
        file_stream.seek(0)
        
        headers = {
            "Content-Disposition": "attachment; filename=SovereignX_Report.pptx"
        }
        return StreamingResponse(
            file_stream,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers=headers
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to generate PowerPoint: {str(e)}")

@router.post("/generate-xlsx")
async def generate_xlsx(request: AgentInvestigateResponse):
    """
    Consumes an existing /agents/investigate response directly and returns a downloadable .xlsx file.
    """
    try:
        import openpyxl
        from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
        
        wb = openpyxl.Workbook()
        
        # Remove default sheet
        default_sheet = wb.active
        wb.remove(default_sheet)
        
        # Format answer and extract unique citations
        formatted_answer, unique_citations = extract_unique_citations(request.answer, request.retrieved_chunks)
        
        # Styling Setup
        header_font = Font(name="Calibri", size=11, bold=True, color="FFFFFF")
        header_fill = PatternFill(start_color="102C57", end_color="102C57", fill_type="solid")
        title_font = Font(name="Calibri", size=14, bold=True, color="102C57")
        bold_font = Font(name="Calibri", size=11, bold=True)
        italic_font = Font(name="Calibri", size=9, italic=True, color="787878")
        alignment_center = Alignment(horizontal="center", vertical="center")
        alignment_left = Alignment(horizontal="left", vertical="center", wrap_text=True)
        thin_side = Side(border_style="thin", color="D3D3D3")
        thin_border = Border(left=thin_side, right=thin_side, top=thin_side, bottom=thin_side)
        
        # SHEET 1: Provenance Audit
        ws_audit = wb.create_sheet(title="Provenance Audit")
        
        # Title Blocks
        ws_audit["A1"] = "SovereignX Grounded Findings Claim-by-Claim Audit"
        ws_audit["A1"].font = title_font
        
        ws_audit["A2"] = f"Query: \"{request.query}\""
        ws_audit["A2"].font = bold_font
        
        ws_audit["A3"] = f"Average Retrieval Relevance Confidence: {round(request.confidence * 100, 2)}%  (Note: represents average similarity of chunks used)"
        ws_audit["A3"].font = italic_font
        
        # Column headers (row 5)
        headers_list = ["Claim Text / Grounded Finding", "Citation Ref", "Source Filename", "Page #", "Source Chunk ID", "Related Tool Verification"]
        for col_num, header in enumerate(headers_list, 1):
            cell = ws_audit.cell(row=5, column=col_num)
            cell.value = header
            cell.font = header_font
            cell.fill = header_fill
            cell.alignment = alignment_center
            cell.border = thin_border
            
        # Parse sentences into claims
        sentences = [s.strip() for s in re.split(r'\.\s+', formatted_answer) if s.strip()]
        
        row_num = 6
        for sentence in sentences:
            if not sentence.endswith('.'):
                sentence += '.'
                
            # Find citation markers inside the sentence
            markers = re.findall(r'\[(\d+)\]', sentence)
            
            # Match with related deterministic tool check summaries (reusing Stage 1 scoped output)
            # The tool result is only mapped to the claim row if the sentence contains the specific numerical value analyzed by the tool.
            related_tool_desc = "N/A"
            for t in request.tool_executions:
                summary = t.get("outputs", {}).get("summary", "")
                
                match_vals = re.search(r'Reading \(([\d.]+).*?SOP limit \(([\d.]+)', summary, re.IGNORECASE)
                if match_vals:
                    reading_val = match_vals.group(1)
                    limit_val = match_vals.group(2)
                    
                    try:
                        r_float = float(reading_val)
                        r_str_1 = f"{r_float}"
                        r_str_2 = f"{int(r_float)}" if r_float.is_integer() else r_str_1
                        
                        l_float = float(limit_val)
                        l_str_1 = f"{l_float}"
                        l_str_2 = f"{int(l_float)}" if l_float.is_integer() else l_str_1
                    except ValueError:
                        r_str_1 = r_str_2 = reading_val
                        l_str_1 = l_str_2 = limit_val
                        
                    # Check if the specific reading value appears in the claim text
                    if r_str_1 in sentence or r_str_2 in sentence:
                        related_tool_desc = summary
                        break
                        
            if not markers:
                ws_audit.cell(row=row_num, column=1, value=sentence).alignment = alignment_left
                ws_audit.cell(row=row_num, column=2, value="N/A").alignment = alignment_center
                ws_audit.cell(row=row_num, column=3, value="N/A").alignment = alignment_left
                ws_audit.cell(row=row_num, column=4, value="N/A").alignment = alignment_center
                ws_audit.cell(row=row_num, column=5, value="N/A").alignment = alignment_center
                ws_audit.cell(row=row_num, column=6, value=related_tool_desc).alignment = alignment_left
                
                for col in range(1, 7):
                    ws_audit.cell(row=row_num, column=col).border = thin_border
                row_num += 1
            else:
                for m in markers:
                    idx = int(m)
                    citation = next((c for c in unique_citations if c["index"] == idx), None)
                    
                    if citation:
                        ws_audit.cell(row=row_num, column=1, value=sentence).alignment = alignment_left
                        ws_audit.cell(row=row_num, column=2, value=f"[{idx}]").alignment = alignment_center
                        ws_audit.cell(row=row_num, column=3, value=citation["filename"]).alignment = alignment_left
                        ws_audit.cell(row=row_num, column=4, value=citation["page"] if citation["page"] is not None else "N/A").alignment = alignment_center
                        ws_audit.cell(row=row_num, column=5, value=citation["chunk_id"]).alignment = alignment_center
                        ws_audit.cell(row=row_num, column=6, value=related_tool_desc).alignment = alignment_left
                        
                        for col in range(1, 7):
                            ws_audit.cell(row=row_num, column=col).border = thin_border
                        row_num += 1
                        
        # Auto-fit columns
        for col in ws_audit.columns:
            max_len = 0
            for cell in col:
                val_str = str(cell.value or '')
                if len(val_str) > max_len:
                    max_len = len(val_str)
            col_letter = openpyxl.utils.get_column_letter(col[0].column)
            ws_audit.column_dimensions[col_letter].width = min(max(max_len + 3, 10), 60)
            
        # SHEET 2: Tool Verifications
        ws_tools = wb.create_sheet(title="Tool Verifications")
        
        ws_tools["A1"] = "Deterministic Tool Verification Logs"
        ws_tools["A1"].font = title_font
        
        ws_tools["A2"] = "Note: The following results are computed deterministically by SovereignX's comparison rules, completely independent of the LLM."
        ws_tools["A2"].font = italic_font
        
        tool_headers = ["Tool Name", "Inputs / Parameters", "Status", "Deterministic Output Summary"]
        for col_num, header in enumerate(tool_headers, 1):
            cell = ws_tools.cell(row=4, column=col_num)
            cell.value = header
            cell.font = header_font
            cell.fill = header_fill
            cell.alignment = alignment_center
            cell.border = thin_border
            
        t_row = 5
        if not request.tool_executions:
            ws_tools.cell(row=t_row, column=1, value="No comparison tools were triggered for this query context.").font = bold_font
            ws_tools.merge_cells(start_row=t_row, start_column=1, end_row=t_row, end_column=4)
        else:
            for t in request.tool_executions:
                ws_tools.cell(row=t_row, column=1, value=t.get("tool_name", "Unknown")).alignment = alignment_left
                
                summary = t.get("outputs", {}).get("summary", "No summary")
                inputs_str = extract_inputs_from_summary(summary)
                if inputs_str == "N/A":
                    args = t.get("arguments", {})
                    if args:
                        inputs_str = ", ".join([f"{k}={v}" for k, v in args.items()])
                
                ws_tools.cell(row=t_row, column=2, value=inputs_str).alignment = alignment_left
                ws_tools.cell(row=t_row, column=3, value=t.get("status", "success")).alignment = alignment_center
                ws_tools.cell(row=t_row, column=4, value=summary).alignment = alignment_left
                
                for col in range(1, 5):
                    ws_tools.cell(row=t_row, column=col).border = thin_border
                t_row += 1
                
        # Auto-fit columns
        for col in ws_tools.columns:
            max_len = 0
            for cell in col:
                val_str = str(cell.value or '')
                if len(val_str) > max_len:
                    max_len = len(val_str)
            col_letter = openpyxl.utils.get_column_letter(col[0].column)
            ws_tools.column_dimensions[col_letter].width = min(max(max_len + 3, 10), 60)
            
        # Save to memory stream
        file_stream = io.BytesIO()
        wb.save(file_stream)
        file_stream.seek(0)
        
        headers = {
            "Content-Disposition": "attachment; filename=SovereignX_Report.xlsx"
        }
        return StreamingResponse(
            file_stream,
            media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            headers=headers
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to generate Excel sheet: {str(e)}")
